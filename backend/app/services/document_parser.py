import logging
from pathlib import Path
from typing import List, Tuple

import fitz
from docx import Document as DocxDocument
from pptx import Presentation

logger = logging.getLogger(__name__)


def parse_pdf(file_path: str) -> Tuple[str, int]:
    doc = fitz.open(file_path)
    pages: List[str] = []
    for page in doc:
        pages.append(page.get_text())
    doc.close()
    return "\n\n".join(pages), len(pages)


def parse_docx(file_path: str) -> Tuple[str, int]:
    doc = DocxDocument(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs), 1


def parse_pptx(file_path: str) -> Tuple[str, int]:
    prs = Presentation(file_path)
    slides: List[str] = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
        slides.append("\n".join(texts))
    return "\n\n".join(slides), len(prs.slides)


def parse_txt(file_path: str) -> Tuple[str, int]:
    text = Path(file_path).read_text(encoding="utf-8", errors="ignore")
    return text, 1


def parse_document(file_path: str, file_type: str) -> Tuple[str, int]:
    parsers = {
        "pdf": parse_pdf,
        "docx": parse_docx,
        "pptx": parse_pptx,
        "txt": parse_txt,
    }
    parser = parsers.get(file_type.lower())
    if parser is None:
        raise ValueError(f"Unsupported file type: {file_type}")
    return parser(file_path)
